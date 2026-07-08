# -*- coding: utf-8 -*-
"""Build an editable PPTX for BUTLER HOUSE investor deck (20 slides)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.text import MSO_AUTO_SIZE

BASE = r"D:\2026\AI\AI创造营\WorkBuddy\2026-07-07-22-57-06\butler-house-deck"
IMG_COVER = os.path.join(BASE, "assets", "Minimalist_luxury_interior__so_2026-07-07T15-03-24.png")
IMG_MID   = os.path.join(BASE, "assets", "Close_up_still_life_of_a_singl_2026-07-07T15-03-54.png")
IMG_CLOSE = os.path.join(BASE, "assets", "Calm_minimalist_horizon__a_sti_2026-07-07T15-03-54.png")

# palette
NAVY   = RGBColor(0x1A,0x3A,0x5C)
NAVY_D = RGBColor(0x0E,0x24,0x38)
CYAN   = RGBColor(0x00,0xD4,0xFF)
INK    = RGBColor(0x10,0x20,0x2F)
SLATE  = RGBColor(0x5A,0x6B,0x7B)
MIST   = RGBColor(0xF4,0xF7,0xFA)
PAPER  = RGBColor(0xFF,0xFF,0xFF)
LINE   = RGBColor(0xE2,0xE8,0xEF)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
DARKTX = RGBColor(0x9F,0xB4,0xC6)
CARDLINE_D = RGBColor(0x3A,0x5A,0x78)

YAHEI = "Microsoft YaHei"
YAHEI_L = "Microsoft YaHei Light"

SW, SH = Inches(13.333), Inches(7.5)
M, CW = 0.7, 11.933
dark_slides = []

def L(x):
    try:
        x.inches
        return x
    except AttributeError:
        return Inches(x)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    if dark:
        s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
        dark_slides.append(s)
    else:
        s.background.fill.solid(); s.background.fill.fore_color.rgb = PAPER
    return s

def add_card(s, l, t, w, h, fill=PAPER, line=LINE, line_w=1.0, radius=0.07):
    l, t, w, h = L(l), L(t), L(w), L(h)
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try: sp.adjustments[0] = radius
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def txt(s, l, t, w, h, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=YAHEI, italic=False, space_after=4, line_spacing=1.15):
    l, t, w, h = L(l), L(t), L(w), L(h)
    box = s.shapes.add_textbox(l, t, w, h); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font; f.color.rgb = color
    return box

def eyebrow(s, text, top=0.5, dark=False):
    txt(s, M, top, CW, 0.4, text, size=12, color=CYAN, bold=True, font=YAHEI)

def title(s, text, top=0.92, size=30, dark=False):
    txt(s, M, top, CW, 0.9, text, size=size, color=(WHITE if dark else NAVY), bold=False, font=YAHEI_L, line_spacing=1.05)

def lead(s, text, top, w=CW, size=14, color=None):
    txt(s, M, top, w, 1.2, text, size=size, color=(color or SLATE), font=YAHEI, line_spacing=1.6)

def bullets(s, l, t, w, h, items, size=13, color=SLATE, dot=CYAN, line_spacing=1.18, space_after=9):
    l, t, w, h = L(l), L(t), L(w), L(h)
    box = s.shapes.add_textbox(l, t, w, h); tf = box.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after); p.line_spacing = line_spacing
        r1 = p.add_run(); r1.text = "•  "; r1.font.size = Pt(size); r1.font.color.rgb = dot; r1.font.name = YAHEI; r1.font.bold = True
        if isinstance(it, tuple):
            rh = p.add_run(); rh.text = it[0]; rh.font.size = Pt(size); rh.font.color.rgb = NAVY; rh.font.name = YAHEI; rh.font.bold = True
            rb = p.add_run(); rb.text = it[1]; rb.font.size = Pt(size); rb.font.color.rgb = color; rb.font.name = YAHEI
        else:
            rb = p.add_run(); rb.text = it; rb.font.size = Pt(size); rb.font.color.rgb = color; rb.font.name = YAHEI
    return box

def kpi_card(s, l, t, w, h, label, value, sub, val_size=24, dark=False):
    fill = NAVY if dark else PAPER
    line = CARDLINE_D if dark else LINE
    sp = add_card(s, l, t, w, h, fill=fill, line=line)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(15); tf.margin_right = Pt(12); tf.margin_top = Pt(13); tf.margin_bottom = Pt(11)
    p1 = tf.paragraphs[0]; p1.space_after = Pt(8)
    r = p1.add_run(); r.text = label; r.font.size = Pt(11); r.font.name = YAHEI
    r.font.color.rgb = (DARKTX if dark else SLATE); r.font.bold = True
    p2 = tf.add_paragraph(); p2.space_after = Pt(6)
    r = p2.add_run(); r.text = value; r.font.size = Pt(val_size); r.font.name = YAHEI_L
    r.font.color.rgb = (WHITE if dark else NAVY)
    p3 = tf.add_paragraph()
    r = p3.add_run(); r.text = sub; r.font.size = Pt(10.5); r.font.name = YAHEI
    r.font.color.rgb = (DARKTX if dark else SLATE)
    return sp

def make_table(s, l, t, w, h, headers, rows, col_w=None, header_fill=NAVY):
    l, t, w, h = L(l), L(t), L(w), L(h)
    nrows = len(rows) + 1; ncols = len(headers)
    gf = s.shapes.add_table(nrows, ncols, l, t, w, h)
    tbl = gf.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Inches(w.inches * cw / total)
    # header
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = header_fill
        c.margin_left = Pt(8); c.margin_right = Pt(6); c.margin_top = Pt(5); c.margin_bottom = Pt(5)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = htext; r.font.size = Pt(11); r.font.bold = True
        r.font.name = YAHEI; r.font.color.rgb = WHITE
    # body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = PAPER
            c.margin_left = Pt(8); c.margin_right = Pt(6); c.margin_top = Pt(4); c.margin_bottom = Pt(4)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val; r.font.size = Pt(10.5); r.font.name = YAHEI
            r.font.color.rgb = (NAVY if j == 0 else INK); r.font.bold = (j == 0)
    # light bottom borders
    return tbl

def bar_chart(s, l, t, w, h, cats, vals, horizontal=False, color=NAVY,
              highlight_idx=None, highlight_color=CYAN, data_labels=False, title=None):
    l, t, w, h = L(l), L(t), L(w), L(h)
    cd = CategoryChartData(); cd.categories = cats; cd.add_series("v", vals)
    ctype = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    g = s.shapes.add_chart(ctype, l, t, w, h, cd); chart = g.chart
    chart.has_legend = False
    plot = chart.plots[0]; plot.gap_width = 55
    ser = chart.series[0]
    if highlight_idx is not None:
        for i, pt in enumerate(ser.points):
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb = (highlight_color if i == highlight_idx else color)
    else:
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = color
    ca = chart.category_axis; va = chart.value_axis
    ca.tick_labels.font.size = Pt(10); ca.tick_labels.font.name = YAHEI; ca.tick_labels.font.color.rgb = SLATE
    va.tick_labels.font.size = Pt(9); va.tick_labels.font.name = YAHEI; va.tick_labels.font.color.rgb = SLATE
    va.has_major_gridlines = False; va.visible = False
    ca.has_major_gridlines = False
    if data_labels:
        plot.has_data_labels = True; dl = plot.data_labels
        dl.number_format = '0'; dl.number_format_is_linked = False
        dl.font.size = Pt(11); dl.font.name = YAHEI; dl.font.color.rgb = NAVY; dl.font.bold = True
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    if title:
        chart.has_title = True; chart.chart_title.text_frame.text = title
        for r in chart.chart_title.text_frame.paragraphs[0].runs:
            r.font.size = Pt(11); r.font.name = YAHEI; r.font.color.rgb = NAVY; r.font.bold = True
    else:
        chart.has_title = False
    return chart

def line_chart(s, l, t, w, h, cats, vals, color=CYAN, data_labels=True, title=None):
    l, t, w, h = L(l), L(t), L(w), L(h)
    cd = CategoryChartData(); cd.categories = cats; cd.add_series("v", vals)
    g = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, l, t, w, h, cd); chart = g.chart
    chart.has_legend = False
    ser = chart.series[0]
    ser.format.line.color.rgb = color; ser.format.line.width = Pt(2.75)
    try:
        ser.marker.format.fill.solid(); ser.marker.format.fill.fore_color.rgb = color
        ser.marker.format.line.color.rgb = color
    except Exception: pass
    ca = chart.category_axis; va = chart.value_axis
    ca.tick_labels.font.size = Pt(10); ca.tick_labels.font.name = YAHEI; ca.tick_labels.font.color.rgb = SLATE
    va.tick_labels.font.size = Pt(9); va.tick_labels.font.name = YAHEI; va.tick_labels.font.color.rgb = SLATE
    va.has_major_gridlines = False; va.visible = False; ca.has_major_gridlines = False
    if data_labels:
        plot = chart.plots[0]; plot.has_data_labels = True; dl = plot.data_labels
        dl.number_format = '0'; dl.number_format_is_linked = False
        dl.font.size = Pt(11); dl.font.name = YAHEI; dl.font.color.rgb = NAVY; dl.font.bold = True
        dl.position = XL_LABEL_POSITION.ABOVE
    if title:
        chart.has_title = True; chart.chart_title.text_frame.text = title
        for r in chart.chart_title.text_frame.paragraphs[0].runs:
            r.font.size = Pt(11); r.font.name = YAHEI; r.font.color.rgb = NAVY; r.font.bold = True
    else:
        chart.has_title = False
    return chart

# ============ SLIDE 1 · COVER ============
s = slide(False)
s.shapes.add_picture(IMG_COVER, 0, 0, width=SW, height=SH)
add_card(s, 0, 0, Inches(6.13), SH, fill=NAVY, line=NAVY, line_w=0, radius=0)
txt(s, 0.7, 2.0, 5.2, 0.4, "BUTLER HOUSE", size=13, color=CYAN, bold=True)
txt(s, 0.7, 2.45, 5.2, 1.0, "巴特勒管家", size=42, color=WHITE, font=YAHEI_L)
txt(s, 0.7, 4.1, 5.0, 1.4, "以终为始 ——\n做中国高端家庭服务的第一品牌。", size=18, color=RGBColor(0xDB,0xE7,0xF1), font=YAHEI_L, line_spacing=1.4)
txt(s, 0.7, 6.25, 5.2, 0.9, "资本级商业计划书 · CAPITAL-GRADE BUSINESS PLAN\nPre-A 轮融资 · 3000 万人民币 · 机密文件 · 仅供投资人查阅", size=11, color=RGBColor(0x8F,0xA6,0xBA), line_spacing=1.5)

# ============ SLIDE 2 · EXEC SUMMARY ============
s = slide(False)
eyebrow(s, "Executive Summary · 执行摘要")
title(s, "五个数字，看懂这门生意")
cw = (CW - 0.9) / 4
xs = [M + i*(cw+0.3) for i in range(4)]
y = 2.15; h = 2.0
kpi_card(s, xs[0], y, cw, h, "赛道定位", "高端家庭\n服务科技", "Premium Household Services Platform", val_size=20)
kpi_card(s, xs[1], y, cw, h, "商业模式", "订阅+佣金\n+增值", "Subscription · Commission · Upsell", val_size=20)
kpi_card(s, xs[2], y, cw, h, "目标退出", "5–7 年 IPO", "港股 / A 股主板，或战略并购", val_size=22)
kpi_card(s, xs[3], y, cw, h, "融资需求", "0.3–5 亿", "Pre-A 3000万 · A 1.2亿 · B 3–5亿", val_size=22)
add_card(s, M, 4.55, CW, 1.7, fill=MIST, line=MIST)
txt(s, M+0.3, 4.75, 3.0, 0.4, "估值逻辑", size=12, color=SLATE, bold=True)
txt(s, M+0.3, 5.15, 4.5, 0.9, "50 – 100 亿", size=34, color=NAVY, font=YAHEI_L)
txt(s, M+5.2, 4.85, 6.4, 1.2, "按 SaaS + 服务平台的混合估值，以退出时点为锚的目标估值区间。", size=14, color=SLATE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.5)

# ============ SLIDE 3 · EXIT ============
s = slide(False)
eyebrow(s, "Chapter 01 · 以终为始")
title(s, "从退出路径，反推每一步战略")
lead(s, "不为融资而融资。所有产品、运营与组织动作，都对齐 5–7 年后的资本退出窗口。", 1.95)
cw = (CW - 0.8) / 3; xs = [M + i*(cw+0.4) for i in range(3)]; y = 3.2; h = 3.2
for i,(tag,head,val,sub) in enumerate([
    ("路径 A · 最优先","港股 / A 股主板 IPO","50–100 亿","Year 5–6 窗口，对标物业与服务龙头 PS 倍数"),
    ("路径 B · 备选","战略并购","40–80 亿","家电 / 物业巨头（美的 · 海尔 · 碧桂园服务）"),
    ("路径 C · 弹性","美股 IPO","8–15 亿$","若政策窗口放开，提供额外退出弹性"),
]):
    sp = add_card(s, xs[i], y, cw, h)
    tf = sp.text_frame; tf.word_wrap = True; tf.margin_left = Pt(16); tf.margin_top = Pt(16); tf.margin_right = Pt(12)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = tag; r.font.size = Pt(11); r.font.bold = True; r.font.name = YAHEI; r.font.color.rgb = CYAN
    p = tf.add_paragraph(); p.space_before = Pt(14); r = p.add_run(); r.text = head; r.font.size = Pt(19); r.font.name = YAHEI_L; r.font.color.rgb = NAVY
    p = tf.add_paragraph(); p.space_before = Pt(8); r = p.add_run(); r.text = val; r.font.size = Pt(26); r.font.name = YAHEI_L; r.font.color.rgb = NAVY
    p = tf.add_paragraph(); p.space_before = Pt(8); r = p.add_run(); r.text = sub; r.font.size = Pt(11.5); r.font.name = YAHEI; r.font.color.rgb = SLATE
    # top accent
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, L(xs[i]), L(y), L(cw), Inches(0.07)); bar.fill.solid(); bar.fill.fore_color.rgb = (CYAN if i==0 else NAVY); bar.line.fill.background()

# ============ SLIDE 4 · VALUATION ============
s = slide(False)
eyebrow(s, "Chapter 01 · 资本市场估值逻辑")
title(s, "对标公司，与我们的估值锚点")
make_table(s, M, 2.2, 6.0, 2.6, ["公司","估值 / 市值","PS"],
           [["碧桂园服务","300 亿+","2–3x"],["好慷在家","50 亿+","3–4x"],
            ["TaskRabbit","被 IKEA 收购","—"],["管家帮","新三板挂牌","1–2x"]], col_w=[1.2,1.4,0.7])
txt(s, 6.9, 2.2, 5.0, 0.4, "巴特勒的估值锚点", size=13, color=NAVY, bold=True)
bullets(s, 6.9, 2.75, 5.0, 3.6, [
    ("短期（1–3 年）：","按 GMV 的 0.5–1x 估值，重增长率与用户质量。"),
    ("中期（3–5 年）：","按收入的 3–5x 估值，重盈利能力与护城河。"),
    ("长期（5 年+）：","按利润的 15–25x 估值，重品牌溢价与垄断地位。"),
], size=13)

# ============ SLIDE 5 · TAM/SAM/SOM ============
s = slide(False)
eyebrow(s, "Chapter 02 · 赛道选择")
title(s, "千亿级市场，尚无绝对龙头")
bar_chart(s, M, 2.2, 7.2, 3.6, ["TAM","SAM","SOM"], [7500,3000,300], horizontal=True, color=NAVY, highlight_idx=0, data_labels=True)
txt(s, 8.0, 2.2, 4.1, 0.4, "市场分层（亿元 / 年）", size=13, color=NAVY, bold=True)
bullets(s, 8.0, 2.75, 4.1, 3.2, [
    "TAM：7500 亿 – 2.5 万亿（高净值家庭服务）",
    "SAM：3000 亿（一线 + 新一线高净值）",
    "SOM：Year 5 达 300 亿 GMV",
    ("市占率：","从 0.3% → 10%"),
], size=12.5)

# ============ SLIDE 6 · PAIN & OPPORTUNITY ============
s = slide(False)
eyebrow(s, "Chapter 02 · 市场痛点与机会")
title(s, "旧市场的问题，即我们的破局点")
txt(s, M, 2.15, 5.8, 0.4, "现有市场的四大问题", size=13, color=NAVY, bold=True)
bullets(s, M, 2.7, 5.8, 3.6, [
    ("供给侧碎片化：","95% 来自个体或小微，标准化低。"),
    ("信任成本高：","缺乏服务人员背景与专业能力机制。"),
    ("服务不可预期：","质量波动大，无售后保障。"),
    ("数字化程度低：","仍停留在电话 / 微信接单。"),
], size=12.5, dot=NAVY)
txt(s, 6.9, 2.15, 5.0, 0.4, "Butler House 的破局点", size=13, color=CYAN, bold=True)
bullets(s, 6.9, 2.7, 5.0, 2.4, [
    "建立「高端管家」职业认证体系（对标 CFA / CPA）。",
    "构建服务 SOP + 数字化智能调度系统。",
    "打造「管家即服务」Butler-as-a-Service 订阅模式。",
], size=12.5)
add_card(s, 6.9, 5.25, 5.0, 1.3, fill=MIST, line=MIST)
txt(s, 7.1, 5.45, 4.6, 1.0, "差异定位：英仕派的服务品质 + 好慷的互联网思维 + 专为新富阶层设计。", size=12.5, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)

# ============ SLIDE 7 · REVENUE MODEL ============
s = slide(False)
eyebrow(s, "Chapter 03 · 商业模式")
title(s, "收入结构：以订阅为锚，多元变现")
bar_chart(s, M, 2.15, 7.2, 2.9, ["会员订阅","服务佣金","增值服务","B 端 SaaS"], [60,25,12,3], horizontal=True, highlight_idx=0, data_labels=True)
txt(s, M+0.1, 5.15, 7.0, 0.4, "目标占比（毛利率）", size=11.5, color=SLATE, bold=True)
cw = (5.0 - 0.8)/3; xs=[6.95 + i*(cw+0.4) for i in range(3)]
for i,(lab,val,sub) in enumerate([
    ("金卡","¥2,999 /月","基础管家 · 12 次 / 月"),
    ("白金卡","¥6,999 /月","专属管家 + 增值服务"),
    ("黑卡","¥19,999 /月","全天候 + 优先权 + 定制"),
]):
    kpi_card(s, xs[i], 2.15, cw, 2.4, lab, val, sub, val_size=18)
txt(s, 6.95, 4.8, 5.0, 1.6, "会员订阅毛利率 85%+，构成收入基本盘；佣金与增值服务提升 ARPU 与黏性。", size=13, color=SLATE, line_spacing=1.6, anchor=MSO_ANCHOR.TOP)

# ============ SLIDE 8 · UNIT ECONOMICS (dark) ============
s = slide(True)
eyebrow(s, "Chapter 03 · 单位经济模型", dark=True)
title(s, "每个管家，都是一台印钞机", dark=True)
lead(s, "单管家月贡献 GMV 1.5–3.5 万，平台收入 3,000–7,000 元，获客与管理成本仅 800–1,200 元 / 月。", 1.95, color=DARKTX)
cw = (CW - 0.9)/4; xs=[M + i*(cw+0.3) for i in range(4)]; y=3.3; h=3.0
kpi_card(s, xs[0], y, cw, h, "单管家月净利贡献", "2,200–5,800 元", "平台口径，扣除获客与管理成本", dark=True)
kpi_card(s, xs[1], y, cw, h, "LTV / CAC", "15–25 x", "行业顶尖水平，留存周期 18 个月", dark=True)
kpi_card(s, xs[2], y, cw, h, "佣金抽成区间", "15–25 %", "基础 15% → 高端定制 25%", dark=True)
kpi_card(s, xs[3], y, cw, h, "订阅毛利", "85 %+", "高毛利业务，构成收入基本盘", dark=True)

# ============ SLIDE 9 · MOAT (image) ============
s = slide(False)
eyebrow(s, "Chapter 04 · 竞争壁垒与护城河")
title(s, "四重护城河，时间越久越深")
s.shapes.add_picture(IMG_MID, L(M), L(2.1), width=Inches(5.1), height=Inches(3.9))
txt(s, 6.4, 2.1, 5.5, 0.4, "四重护城河", size=13, color=NAVY, bold=True)
bullets(s, 6.4, 2.65, 5.6, 2.8, [
    ("品牌壁垒（10 年+）：","「高端管家 = 巴特勒」的心智占位。"),
    ("网络效应（双边）：","家庭端与管家端互相吸引，效率随规模递增。"),
    ("数据壁垒：","高净值家庭服务大数据，AI 需求预测形成飞轮。"),
    ("供应链壁垒：","自建 Butler Academy + 期权绑定，锁定人才供给。"),
], size=12)
add_card(s, 6.4, 5.55, 5.6, 1.05, fill=MIST, line=MIST)
txt(s, 6.6, 5.72, 5.3, 0.8, "核心认证：联合国际管家学院打造中国首个「星级管家」评定标准。", size=12, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

# ============ SLIDE 10 · COMPETITION MATRIX ============
s = slide(False)
eyebrow(s, "Chapter 04 · 竞争策略矩阵")
title(s, "我们站在空白地带")
make_table(s, M, 2.2, CW, 3.4,
    ["维度","好慷在家","管家帮","英仕派","巴特勒（我们）"],
    [["目标客群","中高端白领","大众家庭","顶级富豪","新富 + 高净值"],
     ["价格带 / 年","5K–2 万","1K–1 万","12 万+","3–24 万"],
     ["技术能力","★★★★","★★","★★★","★★★★★"],
     ["品牌调性","年轻化","传统","奢华老派","现代高端"],
     ["资本运作","★★★★","★★★","★★","★★★★★"],
     ["可复制性","★★★★★","★★★","★★","★★★★☆"]],
    col_w=[1.4,1.3,1.1,1.1,1.6])
txt(s, M, 5.8, CW, 0.8, "价格空白（2–10 万 / 年）、人群空白（新富阶层）、科技空白（AI 调度）、认证空白（行业标准）—— 四重空白，皆为我们所填。", size=12.5, color=SLATE, line_spacing=1.4)

# ============ SLIDE 11 · GROWTH ============
s = slide(False)
eyebrow(s, "Chapter 05 · 增长策略 Go-to-Market")
title(s, "三阶段，从 100 户到 50 万户")
lead(s, "GMV 五年增长 160 倍：验证 → 复制 → 规模化。", 1.95)
bar_chart(s, M, 2.35, 7.4, 3.0, ["Y1","Y2","Y3","Y4","Y5"], [0.5,3,10,30,80], horizontal=False, color=NAVY, highlight_idx=4, data_labels=True)
txt(s, M+0.1, 5.45, 7.2, 0.35, "GMV（亿元）", size=11, color=SLATE, bold=True)
cw=(5.0-0.8)/3; xs=[6.95+i*(cw+0.4) for i in range(3)]
for i,(lab,sub) in enumerate([
    ("种子期 · Y1","聚焦上海顶级社区，打磨 SOP，100 户 / 500 万 GMV。"),
    ("扩张期 · Y2–3","复制 4 城，建 Butler Academy，1 万户 / 5 亿 GMV。"),
    ("规模化 · Y4–5","覆盖 15 城，推 SaaS，启动 IPO，10–50 万户 / 50–80 亿。"),
]):
    kpi_card(s, xs[i], 2.35, cw, 2.6, lab, "", sub, val_size=18)

# ============ SLIDE 12 · ACQUISITION ============
s = slide(False)
eyebrow(s, "Chapter 05 · 获客策略")
title(s, "高端家庭的五条精准渠道")
bar_chart(s, M, 2.2, 7.4, 3.7,
    ["私行合作 (¥2,000)","高端社区 (¥1,500)","KOL/圈层 (¥1,000)","异业联盟 (¥3,000)","线上精准 (¥4,000)"],
    [30,25,20,15,10], horizontal=True, color=NAVY, highlight_idx=0, data_labels=True)
txt(s, 8.0, 2.2, 4.0, 0.4, "渠道占比 + CAC", size=13, color=NAVY, bold=True)
bullets(s, 8.0, 2.75, 4.0, 3.2, [
    "私行合作占比最高、CAC 低，是核心获客引擎。",
    "高端社区地推带来高信任度家庭。",
    "KOL / 圈层靠口碑裂变，成本最优。",
    "线上精准仅作补量（占比最低、成本最高）。",
], size=12)

# ============ SLIDE 13 · SUPPLY ============
s = slide(False)
eyebrow(s, "Chapter 05 · 管家供给侧")
title(s, "优质管家，从哪来、怎么长")
bar_chart(s, M, 2.2, 6.6, 3.0, ["高端酒店管家","航空乘务长","海归/留学生","传统家政升级"], [40,25,20,15], horizontal=True, highlight_idx=0, data_labels=True)
txt(s, 7.6, 2.2, 4.4, 0.4, "培养路径（Butler Academy）", size=13, color=CYAN, bold=True)
bullets(s, 7.6, 2.75, 4.5, 2.6, [
    ("初级管家：","3 个月培训 → Butler Certificate"),
    ("中级管家：","6 个月实战 → Butler Professional"),
    ("高级管家：","1 年+ 认证 → Butler Master"),
    ("首席管家：","3 年+ 客户认可"),
], size=12)
add_card(s, 7.6, 5.25, 4.5, 1.25, fill=MIST, line=MIST)
txt(s, 7.8, 5.42, 4.2, 1.0, "年培养 500+ 认证管家；与职业院校合作锁定人才，期权绑定核心人才。", size=12, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.35)

# ============ SLIDE 14 · FINANCIAL ============
s = slide(False)
eyebrow(s, "Chapter 06 · 财务预测")
title(s, "五年财务预测（GMV vs 收入）")
lead(s, "总收入由 0.05 亿增至 77 亿，运营利润率 Year 5 达 67%。", 1.95)
bar_chart(s, M, 2.3, 5.9, 3.0, ["Y1","Y2","Y3","Y4","Y5"], [0.5,3,10,30,80], horizontal=False, color=NAVY, data_labels=True, title="GMV（亿）")
line_chart(s, 6.6, 2.3, 5.3, 3.0, ["Y1","Y2","Y3","Y4","Y5"], [0.15,0.9,3,9,24], title="营业收入（亿）")
cw=(CW-0.9)/4; xs=[M+i*(cw+0.3) for i in range(4)]; y=5.55; h=1.45
for i,(lab,val) in enumerate([("Year 3 毛利率","65%"),("Year 5 毛利率","72%"),("Year 5 运营利润","6 亿"),("Year 5 家庭数","5 万")]):
    kpi_card(s, xs[i], y, cw, h, lab, val, "", val_size=24)

# ============ SLIDE 15 · FUNDING ROADMAP ============
s = slide(False)
eyebrow(s, "Chapter 06 · 融资规划")
title(s, "融资路线图：五步走，通向 IPO")
nodes = [("已完成","天使轮","500 万 · 2500 万","MVP · 种子用户"),
         ("Now","Pre-A 轮","3000 万 · 1.5 亿","上海验证 · 团队建设"),
         ("Year 2","A 轮","1.2 亿 · 6 亿","4 城扩张 · Academy"),
         ("Year 3","B 轮","3–5 亿 · 20 亿","全国 15 城 · SaaS"),
         ("Year 4–5","C / Pre-IPO","8–15 亿 · 50 亿+","IPO 准备 · 并购整合")]
n=len(nodes); gap=CW/(n); y=2.6
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, L(M+gap*0.5), Inches(2.93), Inches(CW-gap), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = LINE; line.line.fill.background()
for i,(yr,name,amt,desc) in enumerate(nodes):
    cx = M + gap*i + gap/2
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, L(cx-Inches(0.11)), Inches(2.82), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = (CYAN if yr=="Now" else NAVY); dot.line.fill.background()
    txt(s, cx-Inches(0.9), y+0.35, Inches(1.8), 0.3, yr, size=11, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, cx-Inches(0.9), y+0.7, Inches(1.8), 0.4, name, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER, font=YAHEI_L)
    txt(s, cx-Inches(0.9), y+1.15, Inches(1.8), 0.9, amt+"\n"+desc, size=10, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.25)
add_card(s, M, 5.5, CW, 1.25, fill=MIST, line=MIST)
txt(s, M+0.3, 5.68, CW-0.6, 1.0, "本轮 Pre-A 条款：融资 3000 万 · 投前估值 1.5 亿 · 稀释 16.67% · 1x 非参与清算优先 · 加权平均反稀释 · 投资方 1 名董事会观察员。", size=12, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)

# ============ SLIDE 16 · CAPITAL EFFICIENCY (dark) ============
s = slide(True)
eyebrow(s, "Chapter 06 · 资本效率指标", dark=True)
title(s, "越跑越快，越省越赚", dark=True)
line_chart(s, M, 2.25, 6.8, 3.5, ["Y1","Y3","Y5"], [8,15,20], title="LTV / CAC 走势")
txt(s, M+0.1, 5.85, 6.6, 0.4, "LTV / CAC（x）", size=11, color=DARKTX, bold=True)
cw=(CW-6.8-0.4-0.6)/2; xs=[7.6+i*(cw+0.3) for i in range(4)]
# 2x2 grid of stat cards on right
positions=[(7.6,2.25),(7.6+cw+0.3,2.25),(7.6,4.0),(7.6+cw+0.3,4.0)]
stats=[("回本周期","12→6 月"),("月度流失率","5→2 %"),("NPS","50→70"),("净利率 Y5","57 %")]
for (px,py),(lab,val) in zip(positions,stats):
    kpi_card(s, px, py, cw, 1.5, lab, val, "", val_size=22, dark=True)

# ============ SLIDE 17 · TEAM ============
s = slide(False)
eyebrow(s, "Chapter 07 · 团队与组织")
title(s, "顶级人才，期权充足")
make_table(s, M, 2.15, CW, 2.5, ["岗位","要求","期权"],
    [["CEO","10 年+ 高端服务 / 平台，从 0 到 1 成功案例","15%"],
     ["COO","曾任高端酒店 / 物业高管，懂服务运营","5%"],
     ["CTO","互联网大厂背景，精通平台架构","5%"],
     ["CMO","高净值营销专家，奢侈品 / 私行背景","3%"],
     ["CFO","IPO 经验，熟悉港股 / A 股流程","2%"]], col_w=[1.2,4.4,0.9])
txt(s, M, 4.9, CW, 0.4, "组织规模演进", size=13, color=NAVY, bold=True)
nodes2=[("Year 1","30 人","验证期 · 管家运营 15"),("Year 3","200 人","扩张期 · 管家运营 120"),("Year 5","800 人","规模化 · 含 Academy 500")]
n=3; gap=CW/n; y=5.45
for i,(yr,num,desc) in enumerate(nodes2):
    cx=M+gap*i+gap/2
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL, L(cx-Inches(0.09)), Inches(5.34), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb=NAVY; dot.line.fill.background()
    txt(s, cx-Inches(1.0), y, Inches(2.0), 0.3, yr, size=11, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, cx-Inches(1.0), y+0.35, Inches(2.0), 0.4, num, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER, font=YAHEI_L)
    txt(s, cx-Inches(1.0), y+0.8, Inches(2.0), 0.6, desc, size=10, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.2)
ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, L(M+gap*0.5), Inches(5.42), Inches(CW-gap), Pt(2)); ln.fill.solid(); ln.fill.fore_color.rgb=LINE; ln.line.fill.background()

# ============ SLIDE 18 · RISK ============
s = slide(False)
eyebrow(s, "Chapter 08 · 风险与应对")
title(s, "每一个风险，都有预案")
make_table(s, M, 2.2, CW, 3.4, ["风险","概率","影响","应对策略"],
    [["服务质量","中","高","SOP + 神秘客机制 + 保险赔付"],
     ["人才流失","高","高","期权绑定 + 竞业协议 + 培养体系"],
     ["监管风险","低","中","合规经营 + 行业标准话语权"],
     ["竞争加剧","中","中","深耕高端 + 建立品牌心智"],
     ["经济下行","中","中","灵活套餐 + 锁定核心客户"]], col_w=[1.3,0.8,0.8,4.2])

# ============ SLIDE 19 · MILESTONES ============
s = slide(False)
eyebrow(s, "Chapter 09 · 里程碑与退出准备")
title(s, "通往挂牌的清晰节奏")
nodes3=[("Year 3","股改","完成 A 轮，建立规范财务体系"),("Year 4","Pre-IPO","引入主权基金 / PE"),
        ("Year 5 H1","审计辅导","审计报告 + 上市辅导"),("Year 5 H2","路演","提交招股书"),("Year 6","挂牌","上市")]
n=5; gap=CW/n; y=2.5
ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, L(M+gap*0.5), Inches(2.83), Inches(CW-gap), Pt(2)); ln.fill.solid(); ln.fill.fore_color.rgb=LINE; ln.line.fill.background()
for i,(yr,name,desc) in enumerate(nodes3):
    cx=M+gap*i+gap/2
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL, L(cx-Inches(0.1)), Inches(2.72), Inches(0.2), Inches(0.2))
    dot.fill.solid(); dot.fill.fore_color.rgb=NAVY; dot.line.fill.background()
    txt(s, cx-Inches(0.95), y+0.35, Inches(1.9), 0.3, yr, size=11, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, cx-Inches(0.95), y+0.7, Inches(1.9), 0.4, name, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER, font=YAHEI_L)
    txt(s, cx-Inches(0.95), y+1.15, Inches(1.9), 0.8, desc, size=10, color=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.25)
txt(s, M, 5.2, CW, 0.4, "五大投资亮点", size=13, color=CYAN, bold=True)
bullets(s, M, 5.7, CW, 1.5, [
    ("赛道优质：","千亿级高端服务，尚无绝对龙头。"),
    ("模式稀缺：","订阅 + 平台撮合，资本效率高。"),
    ("壁垒深厚：","认证 + 数据 + 品牌，护城河清晰。"),
    ("退出明确：","港股 / A 股 IPO 窗口友好。"),
    ("团队待组：","顶级人才期权充足。"),
], size=12, space_after=3)

# ============ SLIDE 20 · CLOSING ============
s = slide(False)
s.shapes.add_picture(IMG_CLOSE, L(6.6), 0, width=Inches(6.733), height=SH)
add_card(s, 0, 0, Inches(6.6), SH, fill=NAVY_D, line=NAVY_D, line_w=0, radius=0)
txt(s, 0.7, 2.0, 0.5, 0.4, "以终为始 · From IPO Backward", size=11, color=CYAN, bold=True)
txt(s, 0.7, 2.6, 5.3, 2.2, "以终为始，\n从 IPO 倒推每一步。\n巴特勒管家，中国高端家庭服务的第一品牌。", size=26, color=WHITE, font=YAHEI_L, line_spacing=1.25)
txt(s, 0.7, 6.0, 5.3, 1.2, "投资人关系 · ir@butlerhouse.com\nPre-A 轮 · 3000 万人民币 · 估值 1.5 亿\n机密文件 · 仅供投资人查阅", size=12, color=RGBColor(0x9F,0xB4,0xC6), line_spacing=1.6)

# page numbers
for i, sl in enumerate(prs.slides, 1):
    is_dark = sl in dark_slides
    num = "%02d / 20" % i
    txt(sl, CW+0.2, SH-Inches(0.45), 0.9, 0.3, num, size=10, color=(DARKTX if is_dark else SLATE), align=PP_ALIGN.RIGHT)

out = os.path.join(BASE, "巴特勒管家-融资商业计划书.pptx")
prs.save(out)
print("SAVED:", out, "slides:", len(prs.slides._sldIdLst))
